from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
GREEN = RGBColor(0x22, 0xC5, 0x5E)       # primary green
DARK_GREEN = RGBColor(0x16, 0xA3, 0x4A)
DARK = RGBColor(0x1A, 0x1A, 0x2E)        # near-black
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF5)
MID_GRAY = RGBColor(0xA1, 0xA1, 0xAA)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)      # amber accent


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if transparency:
        shape.fill.fore_color.theme_color = None
    return shape


def add_text_box(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_slide_header(slide, title, subtitle=None, dark_bg=True):
    bg_color = DARK if dark_bg else WHITE
    text_color = WHITE if dark_bg else DARK
    set_bg(slide, bg_color)
    # accent bar top
    add_rect(slide, 0, 0, 10, 0.06, GREEN)
    add_text_box(slide, title, 0.5, 0.15, 9, 0.7, size=28, bold=True,
                 color=text_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.5, 0.75, 9, 0.4, size=14,
                     color=MID_GRAY if dark_bg else RGBColor(0x71, 0x71, 0x7A),
                     align=PP_ALIGN.LEFT)
    # accent bar bottom
    add_rect(slide, 0, 7.44, 10, 0.06, GREEN)


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # completely blank


# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 10, 0.08, GREEN)
add_rect(slide, 0, 7.42, 10, 0.08, GREEN)

# big green circle decoration
circle = slide.shapes.add_shape(9, Inches(6.8), Inches(-0.5), Inches(4.5), Inches(4.5))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x16, 0x52, 0x34)
circle.line.fill.background()

add_text_box(slide, "🛒  PriceBasket", 0.7, 1.5, 8, 1.0, size=42, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, "India's Quick-Commerce Price Comparison Platform",
             0.7, 2.6, 8, 0.7, size=20, bold=False, color=GREEN, align=PP_ALIGN.LEFT)
add_text_box(slide, "Business Model & Revenue Strategy",
             0.7, 3.35, 8, 0.5, size=16, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

add_text_box(slide, "Nikhil Mathur  |  founder@pricebasket.in  |  May 2025",
             0.7, 6.6, 9, 0.4, size=11, color=MID_GRAY, align=PP_ALIGN.LEFT)


# ── SLIDE 2: What We Do ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "What PriceBasket Does",
                 "Real-time grocery price comparison across 8+ Indian quick-commerce platforms")

cards = [
    ("💰  Find Cheapest",
     "Same 500ml Dettol: ₹89 on Blinkit → ₹72 on JioMart\nPriceBasket finds it in 2 seconds"),
    ("🔔  Price Alerts",
     "Users set a target price.\nGet email the moment any platform drops below it."),
    ("🛍️  Cart Optimizer",
     "Add 10 products — we show the single platform or\nsplit order that saves the most."),
]
x_starts = [0.4, 3.6, 6.8]
for (title, body), x in zip(cards, x_starts):
    add_rect(slide, x, 1.4, 3.0, 3.5, RGBColor(0x25, 0x25, 0x3D))
    add_rect(slide, x, 1.4, 3.0, 0.07, GREEN)
    add_text_box(slide, title, x + 0.15, 1.55, 2.7, 0.55, size=13, bold=True, color=WHITE)
    add_text_box(slide, body, x + 0.15, 2.18, 2.7, 2.6, size=11, color=LIGHT_GRAY)

# platforms row
add_text_box(slide, "Platforms:  Blinkit  ·  Zepto  ·  Swiggy Instamart  ·  BigBasket  ·  Flipkart Minutes  ·  Amazon Now  ·  JioMart  ·  DMart",
             0.4, 5.25, 9.2, 0.5, size=11, color=GREEN, align=PP_ALIGN.CENTER)


# ── SLIDE 3: Target Market ───────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "Target Market", "Four addressable segments across consumers and B2B")

segments = [
    ("120M+", "Urban Grocery Shoppers", "Tier 1 city quick-commerce\nusers saving 15–40%/order"),
    ("40M+",  "Deal-Seekers 25–45 yrs", "Price alerts &\nbest-deal notifications"),
    ("5M+",   "MSME / Small Restaurants", "Bulk price comparison\n& vendor shortlisting"),
    ("8+",    "Platform Advertisers", "Blinkit, Zepto, Amazon + FMCG\nbrands reaching high-intent buyers"),
]
x_positions = [0.3, 2.7, 5.1, 7.5]
for (num, label, desc), x in zip(segments, x_positions):
    add_rect(slide, x, 1.35, 2.2, 3.6, RGBColor(0x25, 0x25, 0x3D))
    add_text_box(slide, num, x + 0.1, 1.5, 2.0, 0.8, size=30, bold=True,
                 color=GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, x + 0.1, 2.3, 2.0, 0.65, size=12, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.1, 3.05, 2.0, 1.5, size=10,
                 color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 4: Revenue Streams Overview ────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "7 Revenue Streams", "Diversified monetisation — from Day 1 to Scale")

streams = [
    ("1", "Affiliate / Referral",    "0.5–3% per order",          "Short Term",   GREEN),
    ("2", "Sponsored Listings",      "CPM ₹150–₹400 / CPC ₹3–₹8", "Medium Term",  ACCENT),
    ("3", "Display Advertising",     "Google AdSense → Direct",    "Short Term",   GREEN),
    ("4", "Newsletter Sponsorship",  "₹5–₹15 per subscriber/mo",   "Short Term",   GREEN),
    ("5", "B2B Data API",            "₹2,999 → ₹9,999/mo plans",   "Medium Term",  ACCENT),
    ("6", "PriceBasket Pro",         "₹99/month subscription",     "Long Term",    RGBColor(0x60, 0xA5, 0xFA)),
    ("7", "Platform Partnership",    "₹15K–₹50K/mo verified badge","Long Term",    RGBColor(0x60, 0xA5, 0xFA)),
]
col_w = [0.35, 2.9, 2.9, 1.7]
headers = ["#", "Stream", "Pricing", "Timeline"]
x_cols = [0.4, 0.85, 3.85, 6.85]
y_start = 1.35

# header row
add_rect(slide, 0.4, y_start, 9.2, 0.38, RGBColor(0x14, 0x14, 0x23))
for header, x in zip(headers, x_cols):
    add_text_box(slide, header, x, y_start + 0.05, 2.7, 0.28,
                 size=10, bold=True, color=GREEN)

for i, (num, name, pricing, timeline, dot_color) in enumerate(streams):
    y = y_start + 0.38 + i * 0.73
    row_bg = RGBColor(0x1E, 0x1E, 0x32) if i % 2 == 0 else RGBColor(0x25, 0x25, 0x3D)
    add_rect(slide, 0.4, y, 9.2, 0.7, row_bg)
    # dot
    dot = slide.shapes.add_shape(9, Inches(x_cols[0]), Inches(y + 0.22),
                                 Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()

    add_text_box(slide, name,    x_cols[1], y + 0.05, 2.7, 0.6, size=11, bold=True,  color=WHITE)
    add_text_box(slide, pricing, x_cols[2], y + 0.05, 2.7, 0.6, size=10, color=MID_GRAY)
    add_text_box(slide, timeline,x_cols[3], y + 0.05, 2.0, 0.6, size=10, color=dot_color)


# ── SLIDE 5: Affiliate + Sponsored (Top 2 streams detail) ────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "Primary Revenue: Affiliate & Sponsored Listings",
                 "Together these drive 80%+ of early revenue")

# Left card
add_rect(slide, 0.4, 1.35, 4.35, 5.7, RGBColor(0x25, 0x25, 0x3D))
add_rect(slide, 0.4, 1.35, 4.35, 0.08, GREEN)
add_text_box(slide, "Affiliate / Referral Commission", 0.55, 1.45, 4.0, 0.55,
             size=14, bold=True, color=WHITE)
lines_a = [
    "• User clicks 'Buy on Blinkit' from PriceBasket",
    "• Referral param passed in URL",
    "• Platform pays 0.5%–3% on completed order",
    "",
    "Formula:",
    "  10,000 orders/mo × ₹400 AOV × 1.5%",
    "  = ₹60,000/month",
    "",
    "Live programs:",
    "  Amazon Associates  ✓",
    "  Flipkart Affiliates  ✓",
    "  Blinkit / Zepto  (direct deals in progress)",
    "",
    "Gross margin: ~100%",
]
add_text_box(slide, "\n".join(lines_a), 0.55, 2.1, 4.0, 4.7, size=10, color=LIGHT_GRAY)

# Right card
add_rect(slide, 5.2, 1.35, 4.35, 5.7, RGBColor(0x25, 0x25, 0x3D))
add_rect(slide, 5.2, 1.35, 4.35, 0.08, ACCENT)
add_text_box(slide, "Sponsored / Featured Listings", 5.35, 1.45, 4.0, 0.55,
             size=14, bold=True, color=WHITE)
lines_s = [
    "Three formats:",
    "  Sponsored Platform  ₹50,000/month",
    "  Sponsored Product   (FMCG brand CPM/CPC)",
    "  Category Sponsor    ₹80,000/month",
    "",
    "Pricing model:",
    "  CPM  ₹150–₹400 per 1,000 impressions",
    "  CPC  ₹3–₹8 per click-through",
    "",
    "Target buyers:",
    "  Blinkit, Zepto, BigBasket (platforms)",
    "  Amul, Nestlé, HUL (FMCG brands)",
    "",
    "Gross margin: ~85%",
]
add_text_box(slide, "\n".join(lines_s), 5.35, 2.1, 4.0, 4.7, size=10, color=LIGHT_GRAY)


# ── SLIDE 6: B2B API + Pro Subscription ──────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "Scale Revenue: Data API & Pro Subscription",
                 "High-margin recurring revenue as data moat compounds")

# API tiers
add_rect(slide, 0.4, 1.35, 4.35, 5.7, RGBColor(0x25, 0x25, 0x3D))
add_rect(slide, 0.4, 1.35, 4.35, 0.08, ACCENT)
add_text_box(slide, "B2B Data API", 0.55, 1.45, 4.0, 0.55,
             size=14, bold=True, color=WHITE)

tiers = [
    ("Starter",    "₹2,999/mo",  "10K API calls · 30-day history",    GREEN),
    ("Growth",     "₹9,999/mo",  "100K calls · 1-year history",        ACCENT),
    ("Enterprise", "Custom",     "Unlimited · real-time webhooks · SLA", RGBColor(0x60, 0xA5, 0xFA)),
]
for i, (tier, price, desc, color) in enumerate(tiers):
    y = 2.1 + i * 1.45
    add_rect(slide, 0.55, y, 4.0, 1.3, RGBColor(0x1A, 0x1A, 0x2E))
    add_text_box(slide, tier,  0.7, y + 0.08, 1.5, 0.4, size=12, bold=True, color=color)
    add_text_box(slide, price, 2.3, y + 0.08, 2.0, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text_box(slide, desc,  0.7, y + 0.55, 3.7, 0.5, size=10, color=MID_GRAY)
add_text_box(slide, "Customers: FMCG brands, market research,\nretail analytics, indie devs",
             0.55, 6.5, 4.0, 0.5, size=10, color=MID_GRAY)

# Pro features
add_rect(slide, 5.2, 1.35, 4.35, 5.7, RGBColor(0x25, 0x25, 0x3D))
add_rect(slide, 5.2, 1.35, 4.35, 0.08, RGBColor(0x60, 0xA5, 0xFA))
add_text_box(slide, "PriceBasket Pro — ₹99/month", 5.35, 1.45, 4.0, 0.55,
             size=14, bold=True, color=WHITE)

features = [
    ("Price alerts",       "3 alerts",          "Unlimited"),
    ("Alert frequency",    "Daily",              "Every 15 min"),
    ("Price history",      "7 days",             "1 year"),
    ("Cart optimizer",     "Basic",              "Advanced"),
    ("CSV export",         "✗",                  "✓"),
    ("Ad-free",            "✗",                  "✓"),
]
# table header
add_rect(slide, 5.35, 2.1, 4.0, 0.38, RGBColor(0x14, 0x14, 0x23))
add_text_box(slide, "Feature",  5.35, 2.13, 1.8, 0.3, size=9, bold=True, color=GREEN)
add_text_box(slide, "Free",     7.2,  2.13, 0.9, 0.3, size=9, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, "Pro",      8.15, 2.13, 1.1, 0.3, size=9, bold=True, color=RGBColor(0x60, 0xA5, 0xFA), align=PP_ALIGN.CENTER)

for i, (feat, free, pro) in enumerate(features):
    y = 2.5 + i * 0.6
    row_bg = RGBColor(0x1E, 0x1E, 0x32) if i % 2 == 0 else RGBColor(0x25, 0x25, 0x3D)
    add_rect(slide, 5.35, y, 4.0, 0.58, row_bg)
    add_text_box(slide, feat, 5.5,  y + 0.1, 1.65, 0.38, size=10, color=WHITE)
    add_text_box(slide, free, 7.2,  y + 0.1, 0.9,  0.38, size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, pro,  8.15, y + 0.1, 1.1,  0.38, size=10, color=RGBColor(0x60, 0xA5, 0xFA), align=PP_ALIGN.CENTER)


# ── SLIDE 7: Revenue Projections ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "Revenue Projections", "Three-stage growth path to ₹50L–₹2Cr/month")

stages = [
    ("Seed\n(Now)",        "10K–50K MAU",   "Affiliate + Ads",              "₹50K–₹2L/mo",   GREEN,                          0.4),
    ("Growth\n6–12 mo",    "100K–500K MAU", "Sponsored listings\n+ Affiliate", "₹5L–₹25L/mo", ACCENT,                         3.55),
    ("Scale\n12–24 mo",    "1M+ MAU",       "All streams active",           "₹50L–₹2Cr/mo",  RGBColor(0x60, 0xA5, 0xFA),    6.7),
]
for (stage, mau, primary, revenue, color, x) in stages:
    add_rect(slide, x, 1.35, 3.0, 5.6, RGBColor(0x25, 0x25, 0x3D))
    add_rect(slide, x, 1.35, 3.0, 0.08, color)
    add_text_box(slide, stage,   x + 0.15, 1.5,  2.7, 0.75, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, "MAU",   x + 0.15, 2.4,  2.7, 0.3,  size=9,  color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, mau,     x + 0.15, 2.7,  2.7, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, primary, x + 0.15, 3.35, 2.7, 0.75, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.2, 4.4, 2.6, 0.04, color)
    add_text_box(slide, "Monthly Revenue", x + 0.15, 4.55, 2.7, 0.3, size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, revenue, x + 0.15, 4.9, 2.7, 0.7, size=17, bold=True, color=color, align=PP_ALIGN.CENTER)

add_text_box(slide, "Break-even at ~50,000 MAU  ·  Positive unit economics from Day 1  ·  No inventory, no COGS",
             0.4, 7.1, 9.2, 0.35, size=11, color=GREEN, align=PP_ALIGN.CENTER)


# ── SLIDE 8: Cost Structure ───────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "Cost Structure", "Lean infra — free tiers carry us to 50K MAU")

items = [
    ("Backend hosting (Render)",    "₹0 (free tier)",         "₹15K–₹60K/mo"),
    ("Frontend hosting (Vercel)",   "₹0–₹1,700/mo",           "₹5K–₹20K/mo"),
    ("PostgreSQL (Neon/Render)",    "₹0–₹3,000/mo",           "₹10K–₹40K/mo"),
    ("Redis (Upstash/Render)",      "₹0–₹2,000/mo",           "₹5K–₹15K/mo"),
    ("Playwright / Scraper infra",  "₹0–₹5,000/mo",           "₹20K–₹80K/mo"),
    ("Email (SMTP/SendGrid)",       "₹0–₹2,000/mo",           "₹5K–₹20K/mo"),
    ("Domain & SSL",                "₹1,000/year",             "₹1,000/year"),
    ("TOTAL",                       "~₹5K–₹15K/mo",           "₹60K–₹2.35L/mo"),
]

add_rect(slide, 0.4, 1.35, 9.2, 0.4, RGBColor(0x14, 0x14, 0x23))
add_text_box(slide, "Cost Item",     0.55, 1.4, 4.5, 0.3, size=10, bold=True, color=GREEN)
add_text_box(slide, "Now (free tier)", 5.15, 1.4, 2.0, 0.3, size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text_box(slide, "At Scale",      7.3, 1.4, 2.2, 0.3, size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

for i, (item, now, scale) in enumerate(items):
    y = 1.75 + i * 0.66
    is_total = item == "TOTAL"
    row_bg = RGBColor(0x16, 0x52, 0x34) if is_total else (RGBColor(0x1E, 0x1E, 0x32) if i % 2 == 0 else RGBColor(0x25, 0x25, 0x3D))
    txt_color = WHITE if is_total else LIGHT_GRAY
    add_rect(slide, 0.4, y, 9.2, 0.64, row_bg)
    add_text_box(slide, item,  0.55, y + 0.1, 4.5, 0.44, size=11, bold=is_total, color=txt_color)
    add_text_box(slide, now,   5.15, y + 0.1, 2.0, 0.44, size=11, bold=is_total, color=txt_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, scale, 7.3,  y + 0.1, 2.2, 0.44, size=11, bold=is_total, color=txt_color, align=PP_ALIGN.CENTER)


# ── SLIDE 9: Competitive Moat ────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "Competitive Moat", "Four compounding advantages that widen over time")

moats = [
    ("1", "Data Network Effect",
     "More scrapers → more accurate prices → more users → platforms want to partner.",
     GREEN),
    ("2", "Price History Database",
     "Unique asset growing daily. Competitors cannot replicate years of historical data overnight.",
     ACCENT),
    ("3", "User Alerts & Engagement Loop",
     "Users with active price alerts return 4–6× per week. Highest retention cohort.",
     RGBColor(0x60, 0xA5, 0xFA)),
    ("4", "India-First Focus",
     "Deep integration with Indian quick-commerce: Blinkit localities, Zepto pincodes, DMart SKUs.",
     RGBColor(0xF4, 0x72, 0x18)),
]
for i, (num, title, desc, color) in enumerate(moats):
    y = 1.4 + i * 1.45
    add_rect(slide, 0.4, y, 9.2, 1.3, RGBColor(0x25, 0x25, 0x3D))
    # number badge
    badge = slide.shapes.add_shape(9, Inches(0.55), Inches(y + 0.3),
                                   Inches(0.6), Inches(0.6))
    badge.fill.solid(); badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_text_box(slide, num, 0.55, y + 0.3, 0.6, 0.6, size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 1.3, y + 0.08, 7.9, 0.45, size=14, bold=True, color=color)
    add_text_box(slide, desc,  1.3, y + 0.58, 7.9, 0.58, size=11, color=LIGHT_GRAY)


# ── SLIDE 10: Key Metrics ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "Key Metrics to Track", "What success looks like at each stage")

metrics = [
    ("MAU",                          "50K → 500K → 1M",    GREEN),
    ("Avg session duration",          "> 3 minutes",         ACCENT),
    ("Price alert activations",       "> 20% of registered", RGBColor(0x60, 0xA5, 0xFA)),
    ("Platform CTR",                  "> 8% of product views",GREEN),
    ("Affiliate conversion rate",     "> 2% of click-throughs",ACCENT),
    ("Returning user rate (D7)",      "> 35%",               RGBColor(0x60, 0xA5, 0xFA)),
]

cols = 3
for i, (metric, target, color) in enumerate(metrics):
    col = i % cols
    row = i // cols
    x = 0.4 + col * 3.2
    y = 1.5 + row * 2.5
    add_rect(slide, x, y, 2.9, 2.1, RGBColor(0x25, 0x25, 0x3D))
    add_rect(slide, x, y, 2.9, 0.06, color)
    add_text_box(slide, target, x + 0.15, y + 0.15, 2.6, 0.75, size=20, bold=True,
                 color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, metric, x + 0.1, y + 0.95, 2.7, 0.9, size=11,
                 color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 11: Go-To-Market ────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_slide_header(slide, "Go-To-Market Strategy", "Organic-first growth → paid & B2B at scale")

gtm = [
    ("🔍", "SEO",
     '"Cheapest [product] India"\n"[product] price Blinkit vs Zepto"\nHigh commercial intent, low competition'),
    ("📱", "Social Media",
     "Instagram Reels + YouTube Shorts\nshowing real savings on everyday groceries\n@pricebasketindia"),
    ("💬", "WhatsApp Community",
     "Weekly best-deals digest to opt-in subscribers\n+91 80058 28390\nZero CAC, high retention"),
    ("🤝", "B2B Outreach",
     "Direct pitch to Blinkit, Zepto, BigBasket\npartner teams for sponsored listing deals"),
    ("🎁", "Referral Program",
     '"Invite a friend → get ₹50 credit"\n(Roadmap)\nViral loop to drive organic MAU growth'),
]
x_positions = [0.3, 2.3, 4.3, 6.3, 8.3]
for (icon, title, body), x in zip(gtm, x_positions):
    add_rect(slide, x, 1.4, 1.7, 5.5, RGBColor(0x25, 0x25, 0x3D))
    add_text_box(slide, icon,  x + 0.1, 1.55, 1.5, 0.65, size=24, align=PP_ALIGN.CENTER, color=WHITE)
    add_text_box(slide, title, x + 0.1, 2.25, 1.5, 0.5, size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, body,  x + 0.1, 2.85, 1.5, 3.8, size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 12: Thank You / CTA ─────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 10, 0.08, GREEN)
add_rect(slide, 0, 7.42, 10, 0.08, GREEN)

circle2 = slide.shapes.add_shape(9, Inches(-1.0), Inches(4.0), Inches(4.5), Inches(4.5))
circle2.fill.solid(); circle2.fill.fore_color.rgb = RGBColor(0x16, 0x52, 0x34)
circle2.line.fill.background()

add_text_box(slide, "Let's Build the Price Intelligence\nLayer for India's Quick Commerce",
             1.0, 1.5, 8.0, 1.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "We're open to partnerships, investor conversations, and platform integrations.",
             1.5, 3.5, 7.0, 0.6, size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)

contacts = [
    ("Nikhil Mathur",        GREEN),
    ("founder@pricebasket.in", WHITE),
    ("pricebasket.in",        GREEN),
    ("+91 80058 28390 (WhatsApp)",  MID_GRAY),
]
y = 4.4
for text, color in contacts:
    add_text_box(slide, text, 2.0, y, 6.0, 0.45, size=13, color=color, align=PP_ALIGN.CENTER)
    y += 0.48


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/nikhilmathur1997/Downloads/Pricebaskettest/docs/PriceBasket_Business_Model.pptx"
prs.save(out)
print("Saved:", out)
